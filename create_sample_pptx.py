#!/usr/bin/env python3
"""
Create a sample PPTX certificate template with placeholders
"""

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except ImportError:
    print("Error: python-pptx not installed. Run: pip install python-pptx")
    exit(1)

def create_sample_certificate():
    # Create presentation
    prs = Presentation()
    
    # Set slide size to standard certificate size (landscape)
    prs.slide_width = Inches(11)
    prs.slide_height = Inches(8.5)
    
    # Add slide
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add border rectangle
    border = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0.5), Inches(0.5),  # Position
        Inches(10), Inches(7.5)    # Size
    )
    border.fill.solid()
    border.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White fill
    border.line.color.rgb = RGBColor(0, 0, 0)  # Black border
    border.line.width = Pt(3)
    
    # Add inner border
    inner_border = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0.75), Inches(0.75),  # Position
        Inches(9.5), Inches(7)       # Size
    )
    inner_border.fill.solid()
    inner_border.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White fill
    inner_border.line.color.rgb = RGBColor(0, 0, 0)  # Black border
    inner_border.line.width = Pt(1)
    
    # Title
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(1.5),    # Position
        Inches(9), Inches(1)       # Size
    )
    title_frame = title_box.text_frame
    title_frame.text = "CERTIFICATE OF COMPLETION"
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_run = title_para.runs[0]
    title_run.font.name = 'Arial'
    title_run.font.size = Pt(28)
    title_run.font.bold = True
    title_run.font.color.rgb = RGBColor(0, 0, 0)
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(
        Inches(1), Inches(2.8),    # Position
        Inches(9), Inches(0.5)     # Size
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "This is to certify that"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.alignment = PP_ALIGN.CENTER
    subtitle_run = subtitle_para.runs[0]
    subtitle_run.font.name = 'Arial'
    subtitle_run.font.size = Pt(16)
    subtitle_run.font.color.rgb = RGBColor(0, 0, 0)
    
    # Name placeholder
    name_box = slide.shapes.add_textbox(
        Inches(1), Inches(3.5),    # Position
        Inches(9), Inches(0.8)     # Size
    )
    name_frame = name_box.text_frame
    name_frame.text = "{{name}}"
    name_para = name_frame.paragraphs[0]
    name_para.alignment = PP_ALIGN.CENTER
    name_run = name_para.runs[0]
    name_run.font.name = 'Arial'
    name_run.font.size = Pt(24)
    name_run.font.bold = True
    name_run.font.color.rgb = RGBColor(0, 0, 139)  # Dark blue
    
    # Achievement text
    achievement_box = slide.shapes.add_textbox(
        Inches(1), Inches(4.5),    # Position
        Inches(9), Inches(0.5)     # Size
    )
    achievement_frame = achievement_box.text_frame
    achievement_frame.text = "has successfully completed the course"
    achievement_para = achievement_frame.paragraphs[0]
    achievement_para.alignment = PP_ALIGN.CENTER
    achievement_run = achievement_para.runs[0]
    achievement_run.font.name = 'Arial'
    achievement_run.font.size = Pt(16)
    achievement_run.font.color.rgb = RGBColor(0, 0, 0)
    
    # Course name placeholder
    course_box = slide.shapes.add_textbox(
        Inches(1), Inches(5.2),    # Position
        Inches(9), Inches(0.6)     # Size
    )
    course_frame = course_box.text_frame
    course_frame.text = "{{course}}"
    course_para = course_frame.paragraphs[0]
    course_para.alignment = PP_ALIGN.CENTER
    course_run = course_para.runs[0]
    course_run.font.name = 'Arial'
    course_run.font.size = Pt(20)
    course_run.font.bold = True
    course_run.font.color.rgb = RGBColor(0, 0, 0)
    
    # Date and instructor
    date_box = slide.shapes.add_textbox(
        Inches(1), Inches(6.5),    # Position
        Inches(4), Inches(0.5)     # Size
    )
    date_frame = date_box.text_frame
    date_frame.text = "Date: {{date}}"
    date_para = date_frame.paragraphs[0]
    date_para.alignment = PP_ALIGN.LEFT
    date_run = date_para.runs[0]
    date_run.font.name = 'Arial'
    date_run.font.size = Pt(14)
    date_run.font.color.rgb = RGBColor(0, 0, 0)
    
    instructor_box = slide.shapes.add_textbox(
        Inches(6), Inches(6.5),    # Position
        Inches(4), Inches(0.5)     # Size
    )
    instructor_frame = instructor_box.text_frame
    instructor_frame.text = "Instructor: {{instructor}}"
    instructor_para = instructor_frame.paragraphs[0]
    instructor_para.alignment = PP_ALIGN.RIGHT
    instructor_run = instructor_para.runs[0]
    instructor_run.font.name = 'Arial'
    instructor_run.font.size = Pt(14)
    instructor_run.font.color.rgb = RGBColor(0, 0, 0)
    
    # Organization
    org_box = slide.shapes.add_textbox(
        Inches(1), Inches(7.2),    # Position
        Inches(9), Inches(0.5)     # Size
    )
    org_frame = org_box.text_frame
    org_frame.text = "{{organization}}"
    org_para = org_frame.paragraphs[0]
    org_para.alignment = PP_ALIGN.CENTER
    org_run = org_para.runs[0]
    org_run.font.name = 'Arial'
    org_run.font.size = Pt(12)
    org_run.font.italic = True
    org_run.font.color.rgb = RGBColor(128, 128, 128)  # Gray
    
    # Save the presentation
    prs.save('certificate-template.pptx')
    print("Sample PPTX certificate template created: certificate-template.pptx")
    print("\nPlaceholders included:")
    print("- {{name}} - Person's name")
    print("- {{course}} - Course name")
    print("- {{date}} - Date")
    print("- {{instructor}} - Instructor name")
    print("- {{organization}} - Organization name")

if __name__ == '__main__':
    create_sample_certificate()
