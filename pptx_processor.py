#!/usr/bin/env python3
"""
PPTX Certificate Generator
Processes PowerPoint templates with placeholders and converts to PDF
"""

import sys
import json
import os
from pathlib import Path
import argparse

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    import comtypes.client
except ImportError as e:
    print(f"Error: Required Python packages not installed: {e}")
    print("Please install: pip install python-pptx comtypes")
    sys.exit(1)

class PPTXProcessor:
    def __init__(self):
        self.supported_placeholders = [
            '{{name}}', '{{NAME}}', '{{Name}}',
            '{{date}}', '{{DATE}}', '{{Date}}',
            '{{course}}', '{{COURSE}}', '{{Course}}',
            '{{instructor}}', '{{INSTRUCTOR}}', '{{Instructor}}',
            '{{organization}}', '{{ORGANIZATION}}', '{{Organization}}'
        ]
    
    def find_and_replace_text(self, presentation, replacements):
        """
        Find and replace placeholder text in all slides and shapes
        """
        replaced_count = 0
        
        for slide in presentation.slides:
            # Process text in shapes
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame'):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            original_text = run.text
                            new_text = original_text
                            
                            # Replace all placeholders
                            for placeholder, value in replacements.items():
                                if placeholder in new_text:
                                    new_text = new_text.replace(placeholder, str(value))
                                    replaced_count += 1
                            
                            if new_text != original_text:
                                run.text = new_text
                
                # Process text in tables if present
                if hasattr(shape, 'table'):
                    for row in shape.table.rows:
                        for cell in row.cells:
                            original_text = cell.text
                            new_text = original_text
                            
                            for placeholder, value in replacements.items():
                                if placeholder in new_text:
                                    new_text = new_text.replace(placeholder, str(value))
                                    replaced_count += 1
                            
                            if new_text != original_text:
                                cell.text = new_text
        
        return replaced_count
    
    def process_template(self, template_path, output_path, replacements):
        """
        Process a PPTX template with replacements
        """
        try:
            # Load the presentation
            presentation = Presentation(template_path)
            
            # Find and replace text
            replaced_count = self.find_and_replace_text(presentation, replacements)
            
            # Save the modified presentation
            presentation.save(output_path)
            
            return {
                'success': True,
                'output_path': output_path,
                'replacements_made': replaced_count
            }
            
        except Exception as e:
            return {
                'success': False,
                'error': str(e)
            }
    
    def convert_to_pdf(self, pptx_path, pdf_path):
        """
        Convert PPTX to PDF using COM automation (Windows only)
        """
        try:
            # Initialize PowerPoint application
            powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
            powerpoint.Visible = 1
            
            # Open the presentation
            presentation = powerpoint.Presentations.Open(os.path.abspath(pptx_path))
            
            # Export as PDF
            presentation.SaveAs(os.path.abspath(pdf_path), 32)  # 32 = PDF format
            
            # Close presentation and quit PowerPoint
            presentation.Close()
            powerpoint.Quit()
            
            return {
                'success': True,
                'pdf_path': pdf_path
            }
            
        except Exception as e:
            return {
                'success': False,
                'error': str(e)
            }
    
    def process_and_convert(self, template_path, output_dir, name, additional_data=None):
        """
        Complete workflow: process template and convert to PDF
        """
        try:
            # Prepare replacements
            replacements = {
                '{{name}}': name,
                '{{NAME}}': name.upper(),
                '{{Name}}': name.title(),
                '{{date}}': additional_data.get('date', ''),
                '{{course}}': additional_data.get('course', ''),
                '{{instructor}}': additional_data.get('instructor', ''),
                '{{organization}}': additional_data.get('organization', '')
            } if additional_data else {
                '{{name}}': name,
                '{{NAME}}': name.upper(),
                '{{Name}}': name.title()
            }
            
            # Create safe filename
            safe_name = "".join(c for c in name if c.isalnum() or c in (' ', '-', '_')).rstrip()
            safe_name = safe_name.replace(' ', '_')
            
            # Process PPTX
            pptx_output = os.path.join(output_dir, f"certificate-{safe_name}.pptx")
            result = self.process_template(template_path, pptx_output, replacements)
            
            if not result['success']:
                return result
            
            # Convert to PDF
            pdf_output = os.path.join(output_dir, f"certificate-{safe_name}.pdf")
            pdf_result = self.convert_to_pdf(pptx_output, pdf_output)
            
            if pdf_result['success']:
                # Clean up intermediate PPTX file
                try:
                    os.remove(pptx_output)
                except:
                    pass
                
                return {
                    'success': True,
                    'name': name,
                    'filename': f"certificate-{safe_name}.pdf",
                    'path': pdf_output,
                    'replacements_made': result['replacements_made']
                }
            else:
                return pdf_result
                
        except Exception as e:
            return {
                'success': False,
                'error': str(e)
            }

    def process_and_convert_with_data(self, template_path, output_dir, data_obj, additional_data=None):
        """
        Complete workflow: process template with data object and convert to PDF
        """
        try:
            # Get name for filename (fallback to 'certificate' if no name)
            name = data_obj.get('name', data_obj.get('Name', 'certificate'))

            # Prepare dynamic replacements based on data object
            replacements = {}

            # Add all data from the object as potential replacements
            for key, value in data_obj.items():
                if value is not None:
                    # Add various case formats
                    replacements[f'{{{{{key}}}}}'] = str(value)
                    replacements[f'{{{{{key.upper()}}}}}'] = str(value).upper()
                    replacements[f'{{{{{key.title()}}}}}'] = str(value).title()
                    replacements[f'{{{{{key.lower()}}}}}'] = str(value).lower()

            # Add additional data if provided
            if additional_data:
                for key, value in additional_data.items():
                    if value and key != 'fieldMappings':  # Skip fieldMappings metadata
                        replacements[f'{{{{{key}}}}}'] = str(value)
                        replacements[f'{{{{{key.upper()}}}}}'] = str(value).upper()
                        replacements[f'{{{{{key.title()}}}}}'] = str(value).title()
                        replacements[f'{{{{{key.lower()}}}}}'] = str(value).lower()

            print(f"Replacements for {name}: {replacements}", file=sys.stderr)

            # Create safe filename
            safe_name = "".join(c for c in name if c.isalnum() or c in (' ', '-', '_')).rstrip()
            safe_name = safe_name.replace(' ', '_')

            # Process PPTX
            pptx_output = os.path.join(output_dir, f"certificate-{safe_name}.pptx")
            result = self.process_template(template_path, pptx_output, replacements)

            if not result['success']:
                return result

            # Convert to PDF
            pdf_output = os.path.join(output_dir, f"certificate-{safe_name}.pdf")
            pdf_result = self.convert_to_pdf(pptx_output, pdf_output)

            if pdf_result['success']:
                # Clean up intermediate PPTX file
                try:
                    os.remove(pptx_output)
                except:
                    pass

                return {
                    'success': True,
                    'name': name,
                    'filename': f"certificate-{safe_name}.pdf",
                    'path': pdf_output,
                    'replacements_made': result['replacements_made'],
                    'data_used': data_obj
                }
            else:
                return pdf_result

        except Exception as e:
            return {
                'success': False,
                'error': str(e)
            }

def main():
    parser = argparse.ArgumentParser(description='Process PPTX certificate templates')
    parser.add_argument('--template', required=True, help='Path to PPTX template')
    parser.add_argument('--output-dir', required=True, help='Output directory')
    parser.add_argument('--names', help='JSON array of names (legacy)')
    parser.add_argument('--data', help='JSON array of data objects')
    parser.add_argument('--additional', help='Additional data as JSON')

    args = parser.parse_args()

    try:
        # Handle both old and new data formats
        if args.data:
            # New format: array of data objects
            data_objects = json.loads(args.data)
            additional_data = json.loads(args.additional) if args.additional else {}
            print(f"Processing {len(data_objects)} certificates with field mappings", file=sys.stderr)
            print(f"Additional data: {additional_data}", file=sys.stderr)
        elif args.names:
            # Legacy format: array of names
            names = json.loads(args.names)
            data_objects = [{'name': name} for name in names]
            additional_data = json.loads(args.additional) if args.additional else {}
            print(f"Processing {len(names)} certificates (legacy format)", file=sys.stderr)
        else:
            raise ValueError("Either --data or --names must be provided")

        processor = PPTXProcessor()
        results = []

        # Ensure output directory exists
        os.makedirs(args.output_dir, exist_ok=True)

        for i, data_obj in enumerate(data_objects):
            print(f"Processing certificate {i+1}: {data_obj}", file=sys.stderr)
            result = processor.process_and_convert_with_data(
                args.template,
                args.output_dir,
                data_obj,
                additional_data
            )
            results.append(result)

        # Output results as JSON
        print(json.dumps({
            'success': True,
            'results': results,
            'total_processed': len(data_objects)
        }))

    except Exception as e:
        print(json.dumps({
            'success': False,
            'error': str(e)
        }))

if __name__ == '__main__':
    main()
